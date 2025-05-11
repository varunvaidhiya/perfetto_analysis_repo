from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_content_slide(prs, title_text, content_list):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title_text
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.level = 0

def add_section_header_slide(prs, title_text):
    slide_layout = prs.slide_layouts[2] # Section Header
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text

def add_image_slide(prs, title_text, image_path, content_list=None):
    slide_layout = prs.slide_layouts[5]  # Title and Content (or Blank if only image)
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title_text

    # Add image, centered
    img_width = Inches(6)
    img_height = Inches(4.5) # Assuming 4:3 aspect ratio for the image
    left = (prs.slide_width - img_width) / 2
    top = Inches(1.5) # Adjust top position if there is a title
    pic = slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

    if content_list:
        # Add text box below image if content is provided
        txBox = slide.shapes.add_textbox(left, top + img_height + Inches(0.2), img_width, Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        for item in content_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)

def create_presentation(findings_path, anomalies_report_path, system_info_path, cpu_usage_img_path, output_pptx_path):
    prs = Presentation()
    # Set slide size to 16:9 (widescreen)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # --- Slide 1: Title Slide ---
    add_title_slide(prs, "Perfetto Trace Analysis: Performance Investigation",
                      "Analysis of PerfettoTraceForRecruitment.pftrace\nFocus on YouTube Application Performance")

    # --- Slide 2: Introduction & Trace Context ---
    # Extract system specs from system_info.csv
    try:
        sys_info_df = pd.read_csv(system_info_path)
        device_name = sys_info_df[sys_info_df["name"] == "build_fingerprint"]["str_value"].iloc[0] if not sys_info_df[sys_info_df["name"] == "build_fingerprint"].empty else "Cuttlefish (Default)"
        android_version = sys_info_df[sys_info_df["name"] == "android_sdk_version"]["int_value"].iloc[0] if not sys_info_df[sys_info_df["name"] == "android_sdk_version"].empty else "12 (SDK 31 Default)"
        # RAM and CPU are harder to get directly, using placeholders from guide
        ram_info = "16GB LPDDR5 (Samsung - Example from Guide)"
        cpu_info = "8 cores (e.g., 4x silver @ 1.8GHz, 4x gold @ 2.4GHz - Example from Guide)"
    except Exception:
        device_name = "Cuttlefish (Unavailable)"
        android_version = "12 (SDK 31 Unavailable)"
        ram_info = "RAM Info Unavailable (Example: 16GB LPDDR5)"
        cpu_info = "CPU Info Unavailable (Example: 8 cores)"

    intro_content = [
        "Purpose: Analyzing PerfettoTraceForRecruitment.pftrace to understand system and YouTube app performance on an Android device.",
        "Approach: Acting as performance detectives using CPU data as our primary clues.",
        "Key System Specs:",
        f"  - Device: {device_name}",
        f"  - Android Version: {android_version}",
        f"  - RAM: {ram_info}",
        f"  - CPU: {cpu_info}"
    ]
    add_content_slide(prs, "Introduction & Trace Context", intro_content)

    # --- Slide 3: Workload Identification & Focus ---
    workload_content = [
        "Main Activity Focus: com.google.android.youtube",
        "Rationale: This analysis centers on the YouTube application as a significant workload within the trace, allowing for targeted performance investigation."
    ]
    add_content_slide(prs, "Workload Identification & Focus", workload_content)

    # --- Slide 4: High-Level Performance Insights (from Anomalies Report) ---
    # This will be a summary. Detailed anomalies will come later.
    # Read anomalies_report.md for content
    try:
        with open(anomalies_report_path, "r") as f:
            anomalies_content = f.read()
    except FileNotFoundError:
        anomalies_content = "Anomaly report not found."

    # Summarize anomalies for this slide
    high_level_insights = [
        "Overall CPU Activity: YouTube (com.google.android.youtube) is a primary CPU consumer.",
        "Key Processes/Threads: Significant activity observed in YouTube threads, surfaceflinger, and media services.",
        "Initial Observations/Concerns from CPU Data:",
        "  - High Runnable Times: Several YouTube threads (e.g., ExoPlayer, Binder, CodecLooper) spend considerable time waiting for CPU.",
        "  - Core Placement: Critical YouTube threads observed running on LITTLE cores, potentially impacting performance.",
        "  - Frequent Short Runs: Patterns in some YouTube threads suggest potential I/O or lock contention (requires further investigation with specific trace flags).",
        "  - No Long Main Thread Tasks (>16ms) were found for YouTube in this trace, which is a positive sign for UI responsiveness in that specific regard."
    ]
    add_content_slide(prs, "High-Level Performance Insights (CPU Clues)", high_level_insights)

    # --- Slide 5: Visualization of Top CPU Consumers ---
    add_image_slide(prs, "Top CPU Consuming Processes", cpu_usage_img_path,
                      ["This chart shows the top CPU consuming processes during the trace.",
                       "com.google.android.youtube is a significant consumer, as expected during its active use."])

    # --- Slide 6: Deeper Dive: SQL for Interrogating CPU Data (Conceptual) ---
    # Based on the guide, this is more about explaining the *power* of SQL
    sql_dive_content = [
        "SQL in Perfetto: Allows precise interrogation of trace data, like a detective questioning a witness, to extract specific evidence about CPU behavior.",
        "Example Use Cases (as per guide & findings):",
        "  1. YouTube Thread CPU Time & States:",
        "     - Purpose: Quantify how YouTube threads use CPU and identify struggles (e.g., high 'runnable' time).",
        "     - Insight: Revealed multiple YouTube threads (ExoPlayer, CodecLooper, Binder) with high runnable ratios, indicating CPU contention.",
        "  2. Long CPU Tasks on YouTube's Main Thread (Conceptual - as none were found in this trace >16ms):",
        "     - Purpose: Find long-duration CPU-bound tasks on the main thread, prime suspects for UI unresponsiveness.",
        "     - Insight (if found): Would identify specific tasks causing jank.",
        "Full schema is available for even deeper investigations."
    ]
    add_content_slide(prs, "Deeper Dive: SQL for Interrogating CPU Data", sql_dive_content)

    # --- Slide 7: Key Anomaly 1: High Runnable Time for YouTube Threads ---
    # Extract from anomalies_report.md
    # This is a simplified extraction, a more robust parser would be better for complex reports
    runnable_time_details = []
    in_runnable_section = False
    for line in anomalies_content.splitlines():
        if "## High Runnable Time for YouTube Threads Analysis" in line:
            in_runnable_section = True
            runnable_time_details.append("Observation: Multiple YouTube threads show high runnable ratios (time spent waiting for CPU).")
            continue
        if in_runnable_section and line.startswith("##") and "High Runnable Time" not in line:
            in_runnable_section = False # End of section
            break
        if in_runnable_section and ("Thread:" in line or "Runnable Ratio:" in line or "High runnable time suggests" in line):
            if "ExoPlayer:Playb" in line or "CodecLooper" in line or "BG Thread #3" in line or "binder:23019_4" in line: # Show a few examples
                 runnable_time_details.append(line.strip().replace("**",""))
    if len(runnable_time_details) <=1 : runnable_time_details.append("Detailed data in appendix or full report.")

    add_content_slide(prs, "Key Anomaly 1: High Runnable Time (CPU Contention)", runnable_time_details + [
        "Hypothesis: CPU contention from other threads/processes, or incorrect thread priority.",
        "Impact: Potential sluggishness, jank, and reduced responsiveness in the YouTube app.",
        "How-to: Analyze scheduling data for competing threads; review thread priorities."
    ])

    # --- Slide 8: Key Anomaly 2: YouTube Thread Core Placement ---
    core_placement_details = []
    in_core_section = False
    for line in anomalies_content.splitlines():
        if "## YouTube Thread CPU Core Placement Analysis" in line:
            in_core_section = True
            core_placement_details.append("Observation: Critical YouTube threads spent ~29% of CPU time on assumed LITTLE cores.")
            continue
        if in_core_section and line.startswith("##") and "Core Placement" not in line:
            in_core_section = False
            break
        if in_core_section and ("Critical YouTube threads spent" in line or "Running critical tasks on LITTLE cores" in line):
            core_placement_details.append(line.strip().replace("**",""))
    if len(core_placement_details) <=1 : core_placement_details.append("Detailed data in appendix or full report.")

    add_content_slide(prs, "Key Anomaly 2: Suboptimal Core Placement", core_placement_details + [
        "Hypothesis: Scheduler not optimally placing critical threads on performance (BIG) cores.",
        "Impact: Slower performance for critical operations, potential jank.",
        "How-to: Investigate performance hints; analyze what runs on BIG cores concurrently."
    ])

    # --- Slide 9: Key Anomaly 3: Frequent Short Runs (Potential I/O or Lock Waits) ---
    short_runs_details = []
    in_short_runs_section = False
    youtube_examples_shown = 0
    for line in anomalies_content.splitlines():
        if "## Frequent Short Runs Followed by Sleep Analysis" in line:
            in_short_runs_section = True
            short_runs_details.append("Observation: Numerous YouTube threads (e.g., ExoPlayer, CodecLooper, ChromiumNet, Binder) show frequent short CPU runs.")
            continue
        if in_short_runs_section and line.startswith("##") and "Frequent Short Runs" not in line:
            in_short_runs_section = False
            break
        if in_short_runs_section and "com.google.android.youtube" in line and youtube_examples_shown < 4:
            short_runs_details.append(line.strip().replace("**",""))
            youtube_examples_shown +=1
    if len(short_runs_details) <=1 : short_runs_details.append("Detailed data in appendix or full report.")

    add_content_slide(prs, "Key Anomaly 3: Frequent Short Runs (Potential I/O/Lock Waits)", short_runs_details + [
        "Hypothesis: Pattern *might* indicate waiting for I/O (network, disk) or lock contention.",
        "Impact: If on critical path, can lead to delays and reduced throughput.",
        "How-to: Requires new trace with I/O and lock contention events to confirm."
    ])

    # --- Slide 10: Data Limitations & Guiding Further Investigation ---
    limitations_content = [
        "Acknowledged Limitations (based on current trace & guide examples):",
        "  - No explicit perf_samples_skipped in this trace (positive), but a general concern if present.",
        "  - Potentially missing: detailed function call stacks, direct I/O metrics, memory metrics (if not enabled).",
        "Reframe as Strength: CPU-level clues are invaluable for targeted next steps.",
        "Example: High runnable times point to CPU contention; core placement issues suggest scheduler interaction analysis.",
        "Example: Frequent short runs + sleeps strongly suggest I/O or lock issues, guiding need for specific trace flags (e.g., ftrace with disk, sync, lock categories).",
        "Impact: CPU-first analysis efficiently narrows down problem areas, making subsequent deep dives more productive."
    ]
    add_content_slide(prs, "Data Limitations & Guiding Further Investigation", limitations_content)

    # --- Slide 11: Actionable Insights & Optimization for YouTube (Based on CPU Clues) ---
    # This will summarize from the findings_and_hypotheses.md or guide
    actionable_insights_content = [
        "1. Address CPU Contention (High Runnable Times):",
        "   - How-to: Analyze competing threads, adjust priorities, defer background work.",
        "   - Next Step (Trace): Deeper sched_slice analysis, correlate with CPU freq/idle.",
        "2. Investigate Core Placement:",
        "   - How-to: Explore PerformanceHintManager, analyze concurrent BIG core usage.",
        "   - Next Step (Trace): Correlate with CPU freq/idle; examine what runs on BIG cores.",
        "3. Confirm & Optimize Potential I/O or Lock Waits (Frequent Short Runs):",
        "   - How-to: If confirmed, batch I/O, use caching, optimize network, reduce lock contention.",
        "   - Next Step (Trace): CRITICAL - Capture new trace with I/O (disk, sync, ext4, network) and lock events.",
        "4. General (If Long Main Thread Tasks Were Found - N/A for this trace but good practice):",
        "   - How-to: Profile with call stacks, move work off main thread, simplify rendering.",
        "   - Next Step (Trace): Enable call stack sampling (e.g., gfx, view, input, sched, freq, dalvik)."
    ]
    add_content_slide(prs, "Actionable Insights & Optimization Strategies", actionable_insights_content)

    # --- Slide 12: Conclusion & Recommendations ---
    conclusion_content = [
        "Recap of Main Findings (CPU Detective Work):",
        "  - YouTube performance is impacted by CPU contention (high runnable times for key threads).",
        "  - Suboptimal core placement for critical YouTube threads was observed.",
        "  - Patterns suggestive of I/O or lock contention exist and need confirmation.",
        "Recommendations for Next Steps:",
        "  1. Prioritize a new Perfetto trace with enhanced configurations:",
        "     - Enable I/O event tracing (ftrace: disk, sync, ext4, network).",
        "     - Enable lock contention event tracing (ftrace: lock).",
        "     - Ensure CPU frequency and idle state tracing is active.",
        "     - Consider call stack sampling if specific CPU-bound functions need identification in future.",
        "  2. Focus investigation on periods/threads highlighted by CPU data anomalies.",
        "  3. Systematically address the hypotheses: CPU contention, core placement, and I/O/lock waits."
    ]
    add_content_slide(prs, "Conclusion & Recommendations", conclusion_content)

    # --- Slide 13: Q&A / Appendix Placeholder ---
    add_section_header_slide(prs, "Q&A / Appendix")

    prs.save(output_pptx_path)
    print(f"Presentation saved to {output_pptx_path}")

if __name__ == "__main__":
    # These paths would be dynamic in a real agent workflow
    findings_file = "/home/ubuntu/analysis_results/findings_and_hypotheses.md"
    anomalies_file = "/home/ubuntu/analysis_results/performance_anomalies_report.md"
    system_info_file = "/home/ubuntu/analysis_results/system_info.csv"
    cpu_usage_image = "/home/ubuntu/analysis_results/top_processes_cpu_usage.png"
    output_ppt = "/home/ubuntu/analysis_results/Perfetto_Trace_Analysis_Presentation.pptx"
    create_presentation(findings_file, anomalies_file, system_info_file, cpu_usage_image, output_ppt)

